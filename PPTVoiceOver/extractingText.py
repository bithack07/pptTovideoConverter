from pptx import Presentation
import os
import moviepy.editor as mp
import comtypes.client
from pdf2image import convert_from_path
import generateVoiceOverScripts as GVS
import conversionToVideo as CTV
import conversionToAudio as CTA
#import pptToImage as PTI

ppt_file = 'C:\\PPTVoiceOver\\pptData\\samplepptx.pptx'
presentation_name = os.path.splitext(os.path.basename(ppt_file))[0]
print(presentation_name)
voiceover_scripts=[]
audio_clips = []  # To store all the generated audio clips
images = [] #to store the images
pdf_out_file_name=""


def init_powerpoint():
   powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
   powerpoint.Visible = 1
   return powerpoint

def ppt_to_pdf(powerpoint,inputFileName, outputFileName):
    if outputFileName[-3:] != 'pdf':
        outputFileName += ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, 32) # formatType = 32 for pdf
    deck.Close()

def convert_files_in_folder(powerpoint, folder):
    files = os.listdir(folder)
    for file in files:
        if file[-3:] == "ppt":
            file_path = f'{folder}\\{file}'
            ppt_to_pdf(file_path, file_path)


def write_to_file(data):
   #print("write_to_file")
   
   filename="voiceOvers/"+presentation_name+".txt"
   with open(filename, 'w') as file:
    for slide_text in data:
        file.write(slide_text + '\n\n')


def pdf_to_img(pdf_file, output_folder):
    # Check if folder exists
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)  # If not, create the folder

    images = convert_from_path(pdf_file)
    image_file_paths = []
    for i, image in enumerate(images):
        fname = os.path.join(output_folder, 'image' + str(i) + '.png')
        image.save(fname, "PNG")
        image_file_paths.append(fname)
        
    return image_file_paths


def extract_text_from_ppt(ppt_file):
   # print("In extract_text_from_ppt")
    prs = Presentation(ppt_file)
    all_text = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = ' '.join(run.text for run in paragraph.runs)
                    slide_text.append(para_text)
            elif shape.shape_type == 19:  # check if it's a table
                table = shape.table
                for row in table.rows:
                    row_text = ' | '.join(cell.text for cell in row.cells)
                    slide_text.append(row_text)
        all_text.append(slide_text)
    return all_text

def refining_texts():
   # print("refining_texts")
    slide_texts = extract_text_from_ppt(ppt_file)
    refined_text=[]
    
    for i, slide in enumerate(slide_texts, 1):
     
        slide_no=f"Slide {i}:"
       # print(f"Slide {i}:")
        cleaned_text = ""
        for text in slide:
            
            if text != "Copyright © 2021, Oracle and/or its affiliates. All rights reserved.":
                cleaned_text += text + '\n'
        if cleaned_text != "":        
         #print(slide_no+'\n'+cleaned_text)
         refined_text.append(GVS.refine_slide_text(slide_no+'\n'+cleaned_text))
    #print(refined_text[7])
    return refined_text   

def generating_voice_over_scripts():
    #print("generating_voice_over_scripts")
    refined_texts=refining_texts()
    for i, slide in enumerate(refined_texts, 1):
     voiceover_scripts.append(GVS.generate_voice_over_scripts(slide))
     write_to_file(voiceover_scripts)
    #print(voiceover_scripts) 

    #To test
    # for i, script in enumerate(voiceover_scripts):
    #  print(f"Voiceover script for slide {i+1}:")
    #  print(script)
    #  print()


def generating_voice_over_scripts_audio(images):
   #print("generating_voice_over_scripts_audio")
   generating_voice_over_scripts()
   audio_clips = CTA.texts_to_audio_files(presentation_name,voiceover_scripts)
   #CTV.combine_audio_files(audio_files, presentation_name.replace(' ', '_') + "/combined.mp3")
  # conversion_ppt_to_images()
   print(presentation_name)
   print(audio_clips)
   print(images)
   CTV.create_audio_video_slides(presentation_name,images,audio_clips)

def conversion_ppt_to_images():
    
    powerpoint = init_powerpoint()
    inputFileName = ppt_file   #e.g "C:\\Users\\User\\Desktop\\file.ppt"
    pdf_out_file_name = "C:\\PPTVoiceOver\\temp_pdf\\example_1.pdf" #e.g "C:\\Users\\User\\Desktop\\output.pdf"
    imagesOutputpath = "C:\\PPTVoiceOver\\Images\\example_1\\"
    ppt_to_pdf(powerpoint,inputFileName, pdf_out_file_name)
       # pdf_file = "your pdf file"  # e.g "C:\\Users\\User\\Desktop\\output.pdf"
    images=pdf_to_img(pdf_out_file_name,imagesOutputpath)
    
    #print(images)
    return images
    

images=conversion_ppt_to_images()
generating_voice_over_scripts_audio(images)





